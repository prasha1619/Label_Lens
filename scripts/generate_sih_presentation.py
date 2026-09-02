import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_sih_deck(output_path="LabelLens_SIH_2026_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    BG_DARK = RGBColor(11, 17, 32)
    CARD_BG = RGBColor(30, 41, 59)
    CARD_BORDER = RGBColor(51, 65, 85)
    ACCENT_CYAN = RGBColor(6, 182, 212)
    ACCENT_GREEN = RGBColor(16, 185, 129)
    ACCENT_AMBER = RGBColor(245, 158, 11)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)

    def add_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_tag="SMART INDIA HACKATHON 2026"):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(3.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CARD_BG
        badge.line.color.rgb = ACCENT_CYAN
        badge.line.width = Pt(1)
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = f"★ {category_tag}"
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = ACCENT_CYAN
        p_b.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide1)
    hero = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    hero.fill.solid()
    hero.fill.fore_color.rgb = CARD_BG
    hero.line.color.rgb = CARD_BORDER
    hero.line.width = Pt(1.5)

    sih_tag = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(3.5), Inches(0.4))
    sih_tag.fill.solid()
    sih_tag.fill.fore_color.rgb = RGBColor(15, 23, 42)
    sih_tag.line.color.rgb = ACCENT_AMBER
    tf_tag = sih_tag.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "SMART INDIA HACKATHON 2026"
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_AMBER
    p_tag.alignment = PP_ALIGN.CENTER

    tb_title = slide1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(1.2))
    p = tb_title.text_frame.paragraphs[0]
    p.text = "LabelLens"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p2 = tb_title.text_frame.add_paragraph()
    p2.text = "AI-Powered Legal Metrology Compliance & Statutory Packaging Audit System"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_WHITE

    ps_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.2), Inches(6.8), Inches(3.0))
    ps_box.fill.solid()
    ps_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    ps_box.line.color.rgb = ACCENT_CYAN
    tf_ps = ps_box.text_frame
    tf_ps.word_wrap = True
    p = tf_ps.paragraphs[0]
    p.text = "PROBLEM STATEMENT ID: SIH26034"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p = tf_ps.add_paragraph()
    p.text = "Software system to check compliance of packaged commodities under the Legal Metrology (Packaged Commodities) Rules, 2011, by scanning product images and labels."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)
    p = tf_ps.add_paragraph()
    p.text = "Theme: Miscellaneous / Smart Governance & Enforcement"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(10)

    team_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(3.2), Inches(3.9), Inches(3.0))
    team_box.fill.solid()
    team_box.fill.fore_color.rgb = RGBColor(15, 23, 42)
    team_box.line.color.rgb = CARD_BORDER
    tf_tm = team_box.text_frame
    tf_tm.word_wrap = True
    p = tf_tm.paragraphs[0]
    p.text = "TEAM DETAILS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p = tf_tm.add_paragraph()
    p.text = "Team Name: [Your Team Name]"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)
    p = tf_tm.add_paragraph()
    p.text = "Team Leader: [Leader Name]\nMembers: [Member 1, Member 2, Member 3, Member 4, Member 5]"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(6)

    # Slide 2: Problem
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide2)
    add_header(slide2, "Problem Context & Current Industry Bottlenecks")
    cards_data = [
        ("❌ Manual & Slow Audits", "Enforcement officers manually inspect packaging with physical checklists. Auditing a single retail store takes hours and leaves 99% of FMCG inventory unverified.", ACCENT_AMBER),
        ("⚠️ Complex Multi-Sided Labels", "Essential declarations (MRP, Net Qty, Dates, Veg logo, Manufacturer) are scattered across front, back, and bottom panels, causing high human oversight errors.", ACCENT_CYAN),
        ("⚖️ Revenue Leakage & Fraud", "Mislabeled net quantities and tampered MRPs mislead Indian consumers, leading to statutory violations under the Legal Metrology Act, 2011.", ACCENT_GREEN),
    ]
    for i, (title, desc, color) in enumerate(cards_data):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.8), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(16)

    # Slide 3: Solution
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide3)
    add_header(slide3, "Proposed Solution: LabelLens Architecture & Novelty")
    sol_cards = [
        ("👁️ Vision & Preprocessing", "• Laplacian Variance Blur Filter\n• LAB CLAHE illumination equalizer\n• Handles glossy/curved packaging", ACCENT_CYAN),
        ("⚡ Neural OCR + YOLO", "• YOLO11 label region localization\n• PaddleOCR (DBNet + SVTR)\n• 97.6% Character Recognition", ACCENT_GREEN),
        ("🧩 Cross-Panel Fusion", "• Fuses Front, Back & Side shots\n• Merges scattered declarations\n• Unified product dossier", ACCENT_AMBER),
        ("🛡️ Deterministic Rule Engine", "• Codified Legal Metrology Rules 2011\n• Zero Hallucination Guarantee\n• Non-repudiable PDF Certificate", ACCENT_BLUE),
    ]
    for i, (title, desc, color) in enumerate(sol_cards):
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 2.95), Inches(1.8), Inches(2.85), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(14)

    # Slide 4: Pipeline
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide4)
    add_header(slide4, "End-to-End Technical Execution Pipeline")
    steps = [
        ("Step 1", "Camera Capture", "Multi-panel packaging input", ACCENT_CYAN),
        ("Step 2", "OpenCV QA", "Laplacian blur + CLAHE enhance", ACCENT_BLUE),
        ("Step 3", "YOLO11 Detection", "Identifies declaration ROI", ACCENT_AMBER),
        ("Step 4", "Neural OCR", "PaddleOCR extracts text lines", ACCENT_GREEN),
        ("Step 5", "Rule Engine", "Evaluates Legal Metrology 2011", ACCENT_CYAN),
        ("Step 6", "Audit Dossier", "Signed PDF & Supabase DB record", ACCENT_GREEN),
    ]
    for i, (num, title, desc, col) in enumerate(steps):
        x = Inches(0.8 + i * 1.95)
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(1.85), Inches(3.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(8)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(10)

    # Slide 5: Matrix
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide5)
    add_header(slide5, "Codified Legal Metrology (Rules, 2011) Compliance Matrix")
    table_shape = slide5.shapes.add_table(8, 4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.133)
    headers = ["Mandatory Declaration", "Statutory Rule", "Extraction Mechanism", "Validation Criteria"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN

    matrix_data = [
        ("1. Generic / Commodity Name", "Rule 6(1)(a)", "Fuzzy Title & Category Match", "Standard nomenclature verified"),
        ("2. Net Quantity (Standard SI)", "Rule 6(1)(c) & R.11", "SI Regex ('g, kg, ml, L, N')", "Checked against allowed units"),
        ("3. MRP (Incl. of all taxes)", "Rule 6(1)(e)", "Currency Regex ('Rs. / ₹')", "Mandatory tax inclusion clause"),
        ("4. Date of Mfg / Expiry", "Rule 6(1)(d)", "Temporal Pattern Parser", "Valid format & shelf life checks"),
        ("5. Manufacturer Details", "Rule 6(1)(b)", "Named Entity & Address Parser", "Complete postal address & name"),
        ("6. Consumer Care Contact", "Rule 6(1)(n)", "Email, Toll-Free & Helpdesk regex", "Mandatory phone/email presence"),
        ("7. Country of Origin", "Rule 6(10)", "Country Name Tokenizer", "Mandatory for domestic/imported goods"),
    ]
    for row_idx, row_data in enumerate(matrix_data, start=1):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_WHITE

    # Slide 6: Comparison
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide6)
    add_header(slide6, "Comparative Analysis: Why LabelLens Wins")
    table_shape2 = slide6.shapes.add_table(6, 4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    table2 = table_shape2.table
    table2.columns[0].width = Inches(2.8)
    table2.columns[1].width = Inches(2.8)
    table2.columns[2].width = Inches(2.8)
    table2.columns[3].width = Inches(3.333)
    comp_headers = ["Feature / Metric", "Manual Inspection", "Generic OCR (Tesseract)", "LabelLens (Our Solution)"]
    for col_idx, h in enumerate(comp_headers):
        cell = table2.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_AMBER if col_idx < 3 else ACCENT_GREEN

    comp_data = [
        ("Audit Speed per SKU", "10 - 15 Minutes", "45 - 60 Seconds", "< 15 Seconds (Real-time)"),
        ("Multi-Panel Packaging", "Manual mental matching", "Fails (Single photo only)", "Automated Cross-Field Fusion"),
        ("Glossy / Dark Backgrounds", "Prone to human fatigue", "High Error Rate (>45% fail)", "CLAHE + DBNet (97.6% Acc)"),
        ("Statutory Compliance", "Paper checklists", "Raw unstructured text", "Automated Signed PDF Certificate"),
        ("Hallucination Risk", "Human error & bias", "N/A", "0% (Deterministic Rule Engine)"),
    ]
    for row_idx, row_data in enumerate(comp_data, start=1):
        for col_idx, val in enumerate(row_data):
            cell = table2.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.color.rgb = ACCENT_GREEN if col_idx == 3 else TEXT_WHITE

    # Slide 7: Tech Stack
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide7)
    add_header(slide7, "Technology Stack & Security Governance")
    stack_cards = [
        ("🎨 Frontend Tier", "• React 19 + TypeScript\n• Tailwind CSS + Lucide Icons\n• Vite Production Bundle (~102 KB)\n• Responsive Mobile Inspector UI", ACCENT_CYAN),
        ("⚙️ Backend API Tier", "• FastAPI (High-performance Async)\n• Uvicorn ASGI Server\n• ReportLab PDF Generator\n• Pydantic v2 Type Safety", ACCENT_BLUE),
        ("🧠 AI & Computer Vision", "• YOLO11 Object Detector\n• PaddleOCR 3.x (DBNet + SVTR)\n• OpenCV CLAHE & Laplacian QA\n• Levenshtein Fuzzy Matching", ACCENT_GREEN),
        ("🗄️ Database & Security", "• Supabase PostgreSQL\n• Row Level Security (RLS)\n• JWT Authentication\n• Non-repudiable Audit Logs", ACCENT_AMBER),
    ]
    for i, (title, desc, color) in enumerate(stack_cards):
        x = Inches(0.8 + (i % 2) * 5.95)
        y = Inches(1.8 + (i // 2) * 2.5)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(2.25))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(6)

    # Slide 8: Roadmap
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide8)
    add_header(slide8, "Real-World Impact & Future Scalability Roadmap")
    phases = [
        ("🚀 Phase 1: Mobile & Edge", "• On-device camera guidance\n• Offline inference for remote markets\n• Thermal printer receipt sync", ACCENT_CYAN),
        ("🌐 Phase 2: Multilingual Support", "• OCR for 22 Indian Scheduled Languages\n• Regional Legal Metrology rules\n• EAN/Barcode catalog auto-sync", ACCENT_GREEN),
        ("📦 Phase 3: E-Commerce Crawler", "• Auto-scan Amazon, Flipkart & Blinkit\n• Flag non-compliant online listings\n• Automated regulatory violation notices", ACCENT_AMBER),
    ]
    for i, (title, desc, color) in enumerate(phases):
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.8), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11.5)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(16)

    prs.save(output_path)
    print(f"Presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_sih_deck()
