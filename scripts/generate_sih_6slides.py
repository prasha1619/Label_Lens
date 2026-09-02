import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_6slide_sih_deck(output_path="LabelLens_SIH2026_Official_6Slides.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    BG_COLOR = RGBColor(248, 250, 252)
    HEADER_BLUE = RGBColor(30, 58, 138)
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(203, 213, 225)
    ACCENT_GREEN = RGBColor(16, 185, 129)
    ACCENT_CYAN = RGBColor(14, 165, 233)
    ACCENT_AMBER = RGBColor(245, 158, 11)
    ACCENT_RED = RGBColor(239, 68, 68)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    TEXT_MAIN = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(71, 85, 105)
    FOOTER_BLUE = RGBColor(2, 132, 199)

    def add_base_template(slide, slide_num, header_title):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        team_badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.4), Inches(1.3), Inches(1.1))
        team_badge.fill.solid()
        team_badge.fill.fore_color.rgb = CARD_BG
        team_badge.line.color.rgb = CARD_BORDER
        tf_t = team_badge.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = "Your\nTeam\nName"
        p_t.font.size = Pt(9)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MUTED
        p_t.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.4), Inches(9.3), Inches(0.9))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = header_title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = HEADER_BLUE
        p.alignment = PP_ALIGN.CENTER

        sih_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.5), Inches(0.4), Inches(1.4), Inches(0.9))
        sih_b.fill.solid()
        sih_b.fill.fore_color.rgb = CARD_BG
        sih_b.line.color.rgb = ACCENT_AMBER
        tf_s = sih_b.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.text = "SMART INDIA\nHACKATHON\n2026"
        p_s.font.size = Pt(8)
        p_s.font.bold = True
        p_s.font.color.rgb = HEADER_BLUE
        p_s.alignment = PP_ALIGN.CENTER

        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.05), Inches(13.333), Inches(0.45))
        footer.fill.solid()
        footer.fill.fore_color.rgb = FOOTER_BLUE
        footer.line.fill.background()

        tf_f = footer.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"@SIH Idea submission- Template                                                                                                                                {slide_num}"
        p_f.font.size = Pt(10)
        p_f.font.bold = True
        p_f.font.color.rgb = RGBColor(255, 255, 255)

    # Slide 1: Title Page
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BG_COLOR
    bg1.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.8))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2026"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = HEADER_BLUE
    p1.alignment = PP_ALIGN.CENTER

    tb_sub = slide1.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(0.5))
    p = tb_sub.text_frame.paragraphs[0]
    p.text = "TITLE PAGE: LabelLens"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.alignment = PP_ALIGN.CENTER

    c1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.1), Inches(7.5), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = CARD_BORDER
    c1.line.width = Pt(1.5)
    tf1 = c1.text_frame
    tf1.word_wrap = True

    entries = [
        ("Problem Statement ID –", "SIH26034"),
        ("Problem Statement Title –", "Software system to check compliance of packaged commodities under the Legal Metrology (Packaged Commodities) Rules, 2011, by scanning product images and labels."),
        ("Theme –", "Miscellaneous / Smart Governance & Regulatory Enforcement"),
        ("PS Category –", "Software"),
        ("Team ID –", "[Your Registered Team ID]"),
        ("Team Name –", "[Your Registered Team Name]"),
    ]

    for idx, (label, val) in enumerate(entries):
        p = tf1.paragraphs[0] if idx == 0 else tf1.add_paragraph()
        p.text = f"{label} "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = HEADER_BLUE
        
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(8)

    c_right = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.8))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = RGBColor(238, 242, 255)
    c_right.line.color.rgb = ACCENT_CYAN
    c_right.line.width = Pt(1.5)
    tf_r = c_right.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "🛡️ LabelLens Platform"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = HEADER_BLUE

    points = [
        "⚡ Automated 7-Point Legal Metrology Rule 2011 Audit",
        "👁️ Computer Vision & YOLO11 Region of Interest Detection",
        "🔤 PaddleOCR Deep Learning Text Recognition (97.6% Acc)",
        "🧩 Multi-Panel Front/Back/Side Cross-Field Fusion",
        "📄 Non-repudiable Signed PDF Compliance Certificate",
    ]
    for pt in points:
        p = tf_r.add_paragraph()
        p.text = pt
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MAIN
        p.space_before = Pt(8)

    # Slide 2: Proposed Solution
    slide2 = prs.slides.add_slide(blank_layout)
    add_base_template(slide2, 2, "LabelLens - Proposed Solution")
    s2_data = [
        ("1. Detailed Explanation", "• Multi-Stage Computer Vision & OCR pipeline designed for physical packaging.\n• Automatically extracts MRP, Net Qty, Dates, Manufacturer & Consumer Care.\n• Cross-Panel Fusion combines front, back, and side angles into one unified compliance dossier.", ACCENT_CYAN),
        ("2. How it Addresses Problem", "• Eliminates 15-minute manual inspector checklists with < 15-second AI scans.\n• Solves glossy, reflective, and curved packaging errors using OpenCV LAB-CLAHE.\n• Provides pixel-accurate bounding box evidence for every declaration.", ACCENT_GREEN),
        ("3. Innovation & Uniqueness", "• 100% Deterministic Legal Rule Engine (Zero LLM hallucination in legal verdicts).\n• Multi-panel evidence aggregation across complex packaging shapes.\n• Instant tamper-evident PDF statutory compliance certificate generation.", ACCENT_AMBER),
    ]
    for i, (title, content, color) in enumerate(s2_data):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.8), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = content
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(14)

    # Slide 3: Technical Approach
    slide3 = prs.slides.add_slide(blank_layout)
    add_base_template(slide3, 3, "TECHNICAL APPROACH & METHODOLOGY")
    t_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.8), Inches(4.9))
    t_card.fill.solid()
    t_card.fill.fore_color.rgb = CARD_BG
    t_card.line.color.rgb = ACCENT_CYAN
    t_card.line.width = Pt(1.5)
    tf_t = t_card.text_frame
    tf_t.word_wrap = True
    p = tf_t.paragraphs[0]
    p.text = "🛠️ Technologies Used"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p_tech = tf_t.add_paragraph()
    p_tech.text = "• AI / Vision: YOLO11, PaddleOCR 3.x (DBNet+SVTR), OpenCV, PIL\n\n• Backend: Python 3.11, FastAPI (Async), Uvicorn, ReportLab\n\n• Frontend: React 19, TypeScript, Tailwind CSS, Vite\n\n• Database: PostgreSQL (Supabase) with Row-Level Security (RLS)"
    p_tech.font.size = Pt(10.5)
    p_tech.font.color.rgb = TEXT_MAIN
    p_tech.space_before = Pt(10)

    steps = [
        ("Step 1", "Camera Capture", "Multi-panel packaging input"),
        ("Step 2", "OpenCV QA", "Laplacian blur check + CLAHE"),
        ("Step 3", "YOLO11 Detection", "Identifies declaration ROI boxes"),
        ("Step 4", "Neural OCR", "PaddleOCR text line extraction"),
        ("Step 5", "Rule Engine", "Evaluates Legal Metrology 2011"),
        ("Step 6", "Audit Dossier", "Signed PDF & Supabase DB record"),
    ]
    for idx, (step_num, title, desc) in enumerate(steps):
        row = idx // 3
        col = idx % 3
        x = Inches(4.8 + col * 2.55)
        y = Inches(1.8 + row * 2.5)
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.45), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_GREEN if idx % 2 == 0 else ACCENT_AMBER
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{step_num}: {title}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = HEADER_BLUE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(6)

    # Slide 4: Feasibility
    slide4 = prs.slides.add_slide(blank_layout)
    add_base_template(slide4, 4, "FEASIBILITY AND VIABILITY")
    f_cards = [
        ("1. Feasibility Analysis", "• Technical Feasibility: Built entirely with open-weights models (YOLO11 + PaddleOCR) running efficiently on standard commodity CPU/GPU hardware.\n\n• Economic Feasibility: Zero per-scan API fees. Free cloud hosting on PostgreSQL (Supabase) + edge client inference.\n\n• Operational Feasibility: Ready for deployment on field officers' existing smartphones & web tablets.", ACCENT_CYAN),
        ("2. Potential Challenges & Risks", "• Low-light, blurred photos in crowded market inspections.\n• Glossy & reflective packaging (bottles, metallic pouches).\n• Non-standard date formats (e.g. 'Use within 12M of Pkd').\n• Missing declarations spread across multiple packaging sides.", ACCENT_RED),
        ("3. Strategies to Overcome", "• Real-time Laplacian Blur Validator prompts instant re-capture.\n• LAB-CLAHE illumination normalization removes glare.\n• Fuzzy Levenshtein regex tokenizers handle diverse syntax.\n• Multi-Panel Cross-Field Fusion Engine aggregates all panels.", ACCENT_GREEN),
    ]
    for i, (title, content, color) in enumerate(f_cards):
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.8), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = content
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(12)

    # Slide 5: Impact & Benefits
    slide5 = prs.slides.add_slide(blank_layout)
    add_base_template(slide5, 5, "IMPACT AND BENEFITS")
    impact_cards = [
        ("🎯 Target Audience Impact", "• Legal Metrology Enforcement Officers: Increases inspection throughput by 50x; eliminates manual paperwork.\n\n• Brand Owners & FMCG Manufacturers: Pre-market packaging validation prevents regulatory seizures & penalties.\n\n• Indian Consumers: Protects against underweight fraud, illegal MRP markups, and obscured expiration dates.", ACCENT_CYAN),
        ("💰 Economic & Governance Benefits", "• Prevents Revenue Leakage: Catches misdeclared tax/MRP and unauthorized surcharges.\n\n• 100% Deterministic Compliance: Zero corruption or subjective human bias in statutory audit reports.\n\n• Non-Repudiable Evidence: Bounding box visual proofs attached to downloadable PDF legal notices.", ACCENT_GREEN),
        ("🌱 Social & Environmental Benefits", "• 100% Paperless Audits: Replaces physical ledger record-keeping with encrypted cloud records.\n\n• Consumer Empowerment: Ensures clear allergen and mandatory dietary declarations (Veg/Non-Veg logos).\n\n• Scalable Nationwide: Deployable across rural district offices and mega-retail hubs alike.", ACCENT_AMBER),
    ]
    for i, (title, content, color) in enumerate(impact_cards):
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.8), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = content
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(12)

    # Slide 6: Research & References
    slide6 = prs.slides.add_slide(blank_layout)
    add_base_template(slide6, 6, "RESEARCH AND REFERENCES")
    ref_cards = [
        ("📜 Legal & Statutory References", "1. Legal Metrology Act, 2009 (Act No. 1 of 2010), Ministry of Consumer Affairs, Food & Public Distribution, Govt of India.\n\n2. Legal Metrology (Packaged Commodities) Rules, 2011 — Rules 6, 9, 11, 18 & 24 (Mandatory Declarations on Pre-Packaged Goods).\n\n3. Department of Consumer Affairs E-Commerce Guidelines (2020 Amendment for Digital Packaging Compliance).", ACCENT_BLUE),
        ("🔬 Computer Vision & AI Research", "1. Ultralytics YOLO11: Real-Time Object Detection and Segmentation for Document Region Localization (2024).\n\n2. PaddleOCR 3.x / PaddleX: DBNet (Real-time Scene Text Detection with Differentiable Binarization) & SVTR (Scene Text Recognition Transformer).\n\n3. OpenCV: Laplacian Variance for Fast Blur Detection & Contrast Limited Adaptive Histogram Equalization (CLAHE).", ACCENT_GREEN),
        ("🌐 Project Repository & Live Demo", "• GitHub Repository: github.com/prasha1619/Label_Lens\n\n• Backend API Docs: Interactive Swagger UI (/docs) & ReDoc (/redoc)\n\n• Production Build: React 19 Client + FastAPI High-Throughput Inference Engine", ACCENT_AMBER),
    ]
    for i, (title, content, color) in enumerate(ref_cards):
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.8), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p_desc = tf.add_paragraph()
        p_desc.text = content
        p_desc.font.size = Pt(9.8)
        p_desc.font.color.rgb = TEXT_MAIN
        p_desc.space_before = Pt(12)

    prs.save(output_path)
    print(f"Official 6-slide SIH deck saved successfully to: {output_path}")

if __name__ == "__main__":
    build_6slide_sih_deck()
